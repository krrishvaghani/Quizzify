from fastapi import FastAPI, File, UploadFile, Depends, HTTPException, status
from fastapi.middleware.cors import CORSMiddleware
from fastapi.security import OAuth2PasswordBearer, OAuth2PasswordRequestForm
from fastapi.responses import StreamingResponse
from motor.motor_asyncio import AsyncIOMotorClient
from pydantic import BaseModel, EmailStr
from typing import List, Optional
from datetime import datetime, timedelta
import os
import jwt
from passlib.context import CryptContext
import io
import csv
from pathlib import Path
import random
import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from dotenv import load_dotenv
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_LEFT

# Load environment variables
load_dotenv()

# File processing libraries
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document
import google.generativeai as genai

app = FastAPI(title="Quizify API")

# CORS Configuration
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:3000", "http://localhost:3001", "http://localhost:3002", "http://localhost:5173"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# MongoDB Configuration
MONGODB_URL = os.getenv("MONGODB_URL", "mongodb://localhost:27017")
DATABASE_NAME = os.getenv("DATABASE_NAME", "quizzify_db")
client = AsyncIOMotorClient(MONGODB_URL)
db = client[DATABASE_NAME]

# Security Configuration
SECRET_KEY = os.getenv("SECRET_KEY", "your-secret-key-change-this-in-production")
ALGORITHM = "HS256"
ACCESS_TOKEN_EXPIRE_MINUTES = 1440  # 24 hours

# Email Configuration
SMTP_SERVER = os.getenv("SMTP_SERVER", "smtp.gmail.com")
SMTP_PORT = int(os.getenv("SMTP_PORT", "587"))
SMTP_USERNAME = os.getenv("SMTP_USERNAME", "")
SMTP_PASSWORD = os.getenv("SMTP_PASSWORD", "")
EMAIL_FROM = os.getenv("EMAIL_FROM", SMTP_USERNAME)

pwd_context = CryptContext(schemes=["bcrypt"], deprecated="auto")
oauth2_scheme = OAuth2PasswordBearer(tokenUrl="token")

# Configure Gemini API (you can also use OpenAI or other models)
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY", "")
if GEMINI_API_KEY:
    genai.configure(api_key=GEMINI_API_KEY)
    print(f"✅ Gemini API configured successfully")
else:
    print("⚠️ Warning: GEMINI_API_KEY not found. AI quiz generation will use fallback method.")

# Models
class UserRegister(BaseModel):
    username: str
    email: EmailStr
    password: str
    full_name: Optional[str] = None

class UserLogin(BaseModel):
    email: EmailStr
    password: str

class Token(BaseModel):
    access_token: str
    token_type: str

class OTPVerification(BaseModel):
    email: EmailStr
    otp: str

class ResendOTP(BaseModel):
    email: EmailStr

class MCQOption(BaseModel):
    text: str
    is_correct: bool

class MCQuestion(BaseModel):
    question: str
    options: List[MCQOption]
    explanation: Optional[str] = None
    tags: Optional[List[str]] = []  # Topics/tags for the question

class QuizGenerate(BaseModel):
    num_questions: int = 5
    difficulty: str = "medium"

class QuizShareSettings(BaseModel):
    visibility: str = "unlisted"  # public, unlisted, password_protected
    password: Optional[str] = None
    allow_anonymous: bool = True

class QuizTimerSettings(BaseModel):
    enabled: bool = False
    timer_type: str = "global"  # global or per_question
    global_duration: int = 1800  # seconds (30 minutes default)
    per_question_duration: int = 30  # seconds per question
    auto_submit: bool = True  # auto-submit when time expires
    show_timer: bool = True  # show countdown timer to students
    
class Quiz(BaseModel):
    id: Optional[str] = None
    title: str
    questions: List[MCQuestion]
    created_by: str
    created_at: datetime
    share_settings: Optional[QuizShareSettings] = None
    timer_settings: Optional[QuizTimerSettings] = None

class RoomSettings(BaseModel):
    enable_timer: bool = False
    timer_duration: int = 60  # seconds per question
    shuffle_questions: bool = False
    shuffle_options: bool = False
    attempts_allowed: int = 1
    show_results_immediately: bool = True

class Room(BaseModel):
    id: Optional[str] = None
    title: str
    description: Optional[str] = None
    quiz_id: str
    host_email: str
    room_code: str
    settings: RoomSettings
    status: str = "waiting"  # waiting, active, completed
    participants: List[str] = []
    max_participants: int = 50
    created_at: datetime
    started_at: Optional[datetime] = None

class CreateRoomRequest(BaseModel):
    quiz_id: str
    title: str
    description: Optional[str] = None
    settings: RoomSettings

class JoinRoomRequest(BaseModel):
    room_code: str

class RoomParticipantScore(BaseModel):
    email: str
    username: str
    score: int
    total_questions: int
    completed_at: Optional[datetime] = None

class UpdateProfile(BaseModel):
    username: str
    full_name: Optional[str] = None

class UpdatePassword(BaseModel):
    current_password: str
    new_password: str

class QuizSubmission(BaseModel):
    quiz_id: str
    student_name: str
    student_email: str
    answers: dict  # {question_index: [selected_option_indices]}
    time_taken: int  # seconds
    time_per_question: Optional[dict] = {}  # {question_index: time_in_seconds}

class QuizAttempt(BaseModel):
    id: Optional[str] = None
    quiz_id: str
    student_name: str
    student_email: str
    answers: dict
    score: int
    total_questions: int
    percentage: float
    time_taken: int
    time_per_question: Optional[dict] = {}  # {question_index: time_in_seconds}
    correct_answers: List[int]  # indices of correct questions
    incorrect_answers: List[int]  # indices of incorrect questions
    unanswered: List[int]  # indices of unanswered questions
    submitted_at: datetime

# Helper Functions
def verify_password(plain_password, hashed_password):
    return pwd_context.verify(plain_password, hashed_password)

def generate_otp():
    """Generate a 6-digit OTP"""
    return str(random.randint(100000, 999999))

async def send_otp_email(email: str, otp: str):
    """Send OTP via email"""
    try:
        # Create message
        message = MIMEMultipart("alternative")
        message["Subject"] = "Verify Your Email - Quizzify"
        message["From"] = EMAIL_FROM
        message["To"] = email

        # Create HTML content
        html = f"""
        <html>
          <body style="font-family: Arial, sans-serif; padding: 20px; background-color: #f4f4f4;">
            <div style="max-width: 600px; margin: 0 auto; background-color: white; padding: 30px; border-radius: 10px; box-shadow: 0 2px 10px rgba(0,0,0,0.1);">
              <div style="text-align: center; margin-bottom: 30px;">
                <h1 style="color: #000; margin-bottom: 10px;">🎓 QUIZZIFY</h1>
                <h2 style="color: #333;">Email Verification</h2>
              </div>
              
              <p style="color: #555; font-size: 16px; line-height: 1.6;">
                Thank you for signing up! Please use the following One-Time Password (OTP) to verify your email address:
              </p>
              
              <div style="background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); padding: 20px; border-radius: 8px; text-align: center; margin: 30px 0;">
                <p style="color: white; font-size: 14px; margin: 0 0 10px 0;">Your OTP Code</p>
                <h1 style="color: white; font-size: 36px; letter-spacing: 8px; margin: 0; font-weight: bold;">{otp}</h1>
              </div>
              
              <p style="color: #555; font-size: 14px; line-height: 1.6;">
                This OTP will expire in <strong>10 minutes</strong>. Please do not share this code with anyone.
              </p>
              
              <p style="color: #999; font-size: 12px; margin-top: 30px; padding-top: 20px; border-top: 1px solid #eee;">
                If you didn't request this verification, please ignore this email.
              </p>
            </div>
          </body>
        </html>
        """

        # Attach HTML content
        part = MIMEText(html, "html")
        message.attach(part)

        # Send email
        if SMTP_USERNAME and SMTP_PASSWORD:
            with smtplib.SMTP(SMTP_SERVER, SMTP_PORT) as server:
                server.starttls()
                server.login(SMTP_USERNAME, SMTP_PASSWORD)
                server.send_message(message)
            print(f"✅ OTP email sent successfully to {email}")
            return True
        else:
            print(f"⚠️ Email not configured. OTP for {email}: {otp}")
            return False
    except Exception as e:
        print(f"❌ Failed to send OTP email: {str(e)}")
        return False

def format_time(seconds):
    """Format time in seconds to readable format"""
    minutes = seconds // 60
    remaining_seconds = seconds % 60
    if minutes > 0:
        return f"{minutes} minute{'s' if minutes != 1 else ''} {remaining_seconds} second{'s' if remaining_seconds != 1 else ''}"
    else:
        return f"{remaining_seconds} second{'s' if remaining_seconds != 1 else ''}"

def get_password_hash(password):
    return pwd_context.hash(password)

def create_access_token(data: dict):
    to_encode = data.copy()
    expire = datetime.utcnow() + timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES)
    to_encode.update({"exp": expire})
    encoded_jwt = jwt.encode(to_encode, SECRET_KEY, algorithm=ALGORITHM)
    return encoded_jwt

async def get_current_user(token: str = Depends(oauth2_scheme)):
    credentials_exception = HTTPException(
        status_code=status.HTTP_401_UNAUTHORIZED,
        detail="Could not validate credentials",
        headers={"WWW-Authenticate": "Bearer"},
    )
    try:
        payload = jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
        email: str = payload.get("sub")
        if email is None:
            raise credentials_exception
    except jwt.PyJWTError:
        raise credentials_exception
    
    user = await db.users.find_one({"email": email})
    if user is None:
        raise credentials_exception
    return user

def extract_text_from_pdf(file_content: bytes) -> dict:
    """Extract text from PDF file with structure analysis"""
    pdf_file = io.BytesIO(file_content)
    pdf_reader = PdfReader(pdf_file)
    
    full_text = ""
    structured_content = {
        "full_text": "",
        "headings": [],
        "important_paragraphs": [],
        "key_terms": []
    }
    
    for page in pdf_reader.pages:
        page_text = page.extract_text()
        full_text += page_text + "\n"
        
        # Identify potential headings (all caps, short lines)
        lines = page_text.split('\n')
        for line in lines:
            line = line.strip()
            if line and len(line) < 100:
                # Potential heading: short, mostly uppercase, or ends with colon
                if line.isupper() or line.endswith(':') or (len(line.split()) <= 5 and line[0].isupper()):
                    structured_content["headings"].append(line)
                # Identify important paragraphs (longer, complete sentences)
                elif len(line) > 50 and '. ' in line:
                    structured_content["important_paragraphs"].append(line)
    
    structured_content["full_text"] = full_text
    return structured_content

def extract_text_from_pptx(file_content: bytes) -> dict:
    """Extract text from PowerPoint file with slide structure"""
    pptx_file = io.BytesIO(file_content)
    prs = Presentation(pptx_file)
    
    structured_content = {
        "full_text": "",
        "headings": [],
        "important_paragraphs": [],
        "key_terms": []
    }
    
    full_text = ""
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = f"\n--- Slide {slide_num} ---\n"
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                
                # Title shapes are usually headings
                if shape.shape_type == 1 or (hasattr(shape, 'is_placeholder') and shape.is_placeholder):
                    structured_content["headings"].append(text)
                    slide_text += f"[HEADING] {text}\n"
                else:
                    # Body text
                    if len(text) > 50:
                        structured_content["important_paragraphs"].append(text)
                    slide_text += f"{text}\n"
        
        full_text += slide_text
    
    structured_content["full_text"] = full_text
    return structured_content

def extract_text_from_docx(file_content: bytes) -> dict:
    """Extract text from Word document with heading analysis"""
    docx_file = io.BytesIO(file_content)
    doc = Document(docx_file)
    
    structured_content = {
        "full_text": "",
        "headings": [],
        "important_paragraphs": [],
        "key_terms": []
    }
    
    full_text = ""
    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        
        full_text += text + "\n"
        
        # Check if it's a heading based on style
        if paragraph.style.name.startswith('Heading'):
            structured_content["headings"].append(text)
        # Bold text often indicates important terms
        elif any(run.bold for run in paragraph.runs if run.text.strip()):
            bold_text = ' '.join(run.text for run in paragraph.runs if run.bold and run.text.strip())
            if bold_text and bold_text not in structured_content["key_terms"]:
                structured_content["key_terms"].append(bold_text)
        # Regular paragraphs with substantial content
        elif len(text) > 50 and '. ' in text:
            structured_content["important_paragraphs"].append(text)
    
    structured_content["full_text"] = full_text
    return structured_content

def extract_key_concepts(structured_content: dict, max_concepts: int = 20) -> List[str]:
    """Extract key concepts from structured content"""
    concepts = []
    
    # Prioritize headings as main concepts
    concepts.extend(structured_content.get("headings", [])[:max_concepts // 2])
    
    # Extract important terms from paragraphs
    important_paras = structured_content.get("important_paragraphs", [])
    for para in important_paras[:10]:
        # Extract capitalized terms (likely proper nouns or important concepts)
        words = para.split()
        for i, word in enumerate(words):
            if len(word) > 3 and word[0].isupper() and not word.isupper():
                # Check if it's not at the start of sentence
                if i > 0 or (i == 0 and len(concepts) < max_concepts):
                    concept = word.strip('.,;:!?()')
                    if concept and concept not in concepts:
                        concepts.append(concept)
                        if len(concepts) >= max_concepts:
                            break
        if len(concepts) >= max_concepts:
            break
    
    # Add explicit key terms if available
    key_terms = structured_content.get("key_terms", [])
    for term in key_terms:
        if term not in concepts and len(concepts) < max_concepts:
            concepts.append(term)
    
    return concepts[:max_concepts]

def analyze_content_topics(structured_content: dict) -> str:
    """Analyze content and create a focused summary for AI"""
    headings = structured_content.get("headings", [])
    important_paras = structured_content.get("important_paragraphs", [])
    key_terms = structured_content.get("key_terms", [])
    
    # Build a focused summary emphasizing important content
    summary = ""
    
    if headings:
        summary += "MAIN TOPICS:\n"
        for i, heading in enumerate(headings[:10], 1):
            summary += f"{i}. {heading}\n"
        summary += "\n"
    
    if key_terms:
        summary += "KEY TERMS: " + ", ".join(key_terms[:15]) + "\n\n"
    
    summary += "IMPORTANT CONTENT:\n"
    # Include important paragraphs with context
    for i, para in enumerate(important_paras[:15], 1):
        # Limit paragraph length for efficiency
        para_text = para[:300] + "..." if len(para) > 300 else para
        summary += f"[{i}] {para_text}\n\n"
    
    # Add a portion of full text for context
    full_text = structured_content.get("full_text", "")
    if len(summary) < 3000:
        remaining_chars = 6000 - len(summary)
        summary += "\nADDITIONAL CONTEXT:\n" + full_text[:remaining_chars]
    
    return summary

async def generate_mcqs_with_ai(structured_content: dict, num_questions: int, difficulty: str) -> List[dict]:
    """Generate MCQs using AI with enhanced content analysis"""
    
    # Extract key concepts for better question generation
    key_concepts = extract_key_concepts(structured_content, max_concepts=num_questions * 2)
    
    # Create focused content summary
    analyzed_content = analyze_content_topics(structured_content)
    
    # Map difficulty to specific instructions
    difficulty_instructions = {
        "easy": """
        - Focus on DEFINITIONS and BASIC FACTS from the content
        - Test RECALL and RECOGNITION of key terms and concepts
        - Use direct language and straightforward questions
        - Example: "What is the definition of [term]?" or "Which of the following describes [concept]?"
        """,
        "medium": """
        - Test COMPREHENSION and UNDERSTANDING of relationships between concepts
        - Require APPLICATION of knowledge to new situations
        - Ask about CAUSES, EFFECTS, and COMPARISONS
        - Example: "How does [concept A] relate to [concept B]?" or "What is the primary purpose of [process]?"
        """,
        "hard": """
        - Require ANALYSIS, SYNTHESIS, and EVALUATION of information
        - Test CRITICAL THINKING and ability to draw conclusions
        - Ask about IMPLICATIONS, PREDICTIONS, and COMPLEX RELATIONSHIPS
        - Example: "Based on the principles discussed, what would happen if [scenario]?" or "Evaluate the significance of [concept]"
        """
    }
    
    prompt = f"""You are an EXPERT EXAM-QUESTION GENERATOR with specialized expertise in creating high-quality, meaningful multiple-choice questions for academic examinations.

🎯 CRITICAL MISSION: Generate EXACTLY {num_questions} EXAM-LEVEL multiple-choice questions strictly based on IMPORTANT concepts from the provided document. Every question must be meaningful, useful, and clearly understandable with realistic options.

📚 DOCUMENT CONTENT TO ANALYZE:
{analyzed_content}

🔑 KEY CONCEPTS IDENTIFIED FROM DOCUMENT:
{chr(10).join(f"• {concept}" for concept in key_concepts[:num_questions])}

DIFFICULTY LEVEL: {difficulty.upper()}
{difficulty_instructions.get(difficulty.lower(), difficulty_instructions["medium"])}

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
⚠️ STRICT MANDATORY RULES - FOLLOW WITHOUT EXCEPTION:
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

1. ✅ SOURCE RESTRICTION:
   - Generate questions ONLY from important concepts in the document
   - DO NOT generate random or irrelevant questions
   - Use ONLY information present inside the document
   - NEVER add external knowledge or make assumptions

2. ✅ QUESTION QUALITY:
   - Focus on MAIN TOPICS, KEY CONCEPTS, and CENTRAL IDEAS
   - NO trivial, vague, or superficial questions
   - Each question must be EXAM-LEVEL and test real understanding
   - Questions must be CLEAR, SPECIFIC, and UNAMBIGUOUS
   - DO NOT repeat questions or ask the same concept twice

3. ✅ ANSWER OPTIONS (CRITICAL):
   - Every option must be MEANINGFUL, LOGICAL, and CLOSELY RELATED to the question
   - All 4 options should be REALISTIC and about the SAME TOPIC
   - 1 option is CORRECT based on document
   - 3 options are INCORRECT but PLAUSIBLE (meaningful distractors)
   - Wrong options should be:
     * Related concepts from the document
     * Common misconceptions
     * Partial truths or incomplete answers
     * Similar but incorrect variations
   - AVOID:
     * "None of the above"
     * "All of the above"
     * Obviously wrong options
     * Unrelated or random options
     * Combination options like "Both A and B"

4. ✅ OPTION CONSISTENCY:
   - All options must have similar LENGTH (±20 characters)
   - All options must have similar COMPLEXITY
   - All options must be grammatically parallel
   - All options must be written in the same style

5. ✅ CONCEPT COVERAGE:
   - Each question covers a DIFFERENT important concept
   - Questions should span across multiple sections
   - Prioritize concepts from HEADINGS and KEY TERMS
   - Focus on concepts that demonstrate MASTERY

6. ✅ CLARITY & PROFESSIONALISM:
   - Use professional academic language
   - Questions should be self-contained (no ambiguity)
   - Use terminology from the document
   - Be precise and specific

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
📝 REQUIRED JSON OUTPUT FORMAT:
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

Return ONLY valid JSON array (no markdown, no code blocks, no extra text):

[
  {{
    "question": "What is the primary mechanism by which [specific concept from document] functions?",
    "options": [
      {{"text": "Through process A as described in the document involving specific steps", "is_correct": true}},
      {{"text": "Through process B which is a related but different mechanism", "is_correct": false}},
      {{"text": "Through process C which is mentioned but applies to a different context", "is_correct": false}},
      {{"text": "Through process D which represents a common misconception", "is_correct": false}}
    ],
    "explanation": "The correct answer is A. According to the document, [specific quote or reference from document]. This mechanism is fundamental because [why it matters]."
  }},
  {{
    "question": "According to the document, which statement accurately describes [important concept]?",
    "options": [
      {{"text": "Related but incorrect statement mixing concepts from document", "is_correct": false}},
      {{"text": "Accurate statement directly from document content", "is_correct": true}},
      {{"text": "Partially correct statement missing key details", "is_correct": false}},
      {{"text": "Opposite of what document states but plausible", "is_correct": false}}
    ],
    "explanation": "Option B is correct. The document explicitly states that [specific information]. This is a key concept because [relevance]."
  }}
]

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
⚠️ FINAL CHECKLIST BEFORE SUBMITTING:
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

✅ Generated EXACTLY {num_questions} questions
✅ Each question is from IMPORTANT concepts in document
✅ NO random or irrelevant questions
✅ ALL options are meaningful, logical, and related to question topic
✅ NO "None of the above" or "All of the above"
✅ NO repeated questions or concepts
✅ All options similar in length and complexity
✅ Every question has clear, specific wording
✅ Used ONLY information from the document
✅ Each question tests a DIFFERENT concept
✅ Returned ONLY JSON array (no markdown)

CRITICAL: If you cannot create a high-quality question following ALL rules, skip it and create a different one. Quality over quantity - but you MUST generate exactly {num_questions} questions.
"""
    
    try:
        if GEMINI_API_KEY:
            print(f"🤖 Using Google Gemini AI to generate {num_questions} {difficulty} questions...")
            print(f"📊 Identified {len(key_concepts)} key concepts for question generation")
            
            # Use Gemini Pro for high-quality question generation
            model = genai.GenerativeModel('gemini-pro')
            
            # Optimized generation parameters for educational content
            generation_config = {
                "temperature": 0.3,  # Low temperature for precise, factual questions
                "top_p": 0.85,       # Focus on high-probability outputs
                "top_k": 30,         # Limit to most relevant tokens
                "max_output_tokens": 8192,
                "candidate_count": 1,
            }
            
            response = model.generate_content(
                prompt,
                generation_config=generation_config
            )
            
            # Parse the response and extract JSON
            import json
            import re
            response_text = response.text.strip()
            
            # Remove markdown code blocks if present
            response_text = re.sub(r'^```json\s*', '', response_text)
            response_text = re.sub(r'^```\s*', '', response_text)
            response_text = re.sub(r'\s*```$', '', response_text)
            response_text = response_text.strip()
            
            # Try to find JSON in the response
            json_match = re.search(r'\[[\s\S]*\]', response_text)
            if json_match:
                json_str = json_match.group()
                questions = json.loads(json_str)
                
                print(f"✅ AI generated {len(questions)} questions successfully")
                
                # Validate question structure and quality with scoring
                valid_questions = []
                quality_scores = []
                
                for idx, q in enumerate(questions, 1):
                    if (isinstance(q, dict) and 
                        "question" in q and 
                        "options" in q and 
                        isinstance(q["options"], list) and
                        len(q["options"]) == 4):
                        
                        # Ensure exactly one correct answer
                        correct_count = sum(1 for opt in q["options"] if opt.get("is_correct", False))
                        if correct_count == 1:
                            # Additional quality checks
                            question_text = q["question"].strip()
                            if len(question_text) > 15:  # Reasonable question length
                                # Check if options are not empty and substantial
                                options_valid = all(
                                    opt.get("text", "").strip() and len(opt.get("text", "").strip()) > 5 
                                    for opt in q["options"]
                                )
                                
                                if options_valid:
                                    # Strict quality validation
                                    quality_score = 0
                                    quality_issues = []
                                    
                                    # Check 1: Has explanation (+1)
                                    if q.get("explanation", "").strip() and len(q.get("explanation", "")) > 20:
                                        quality_score += 1
                                    else:
                                        quality_issues.append("weak explanation")
                                    
                                    # Check 2: Question has key concept from document (+1)
                                    if any(concept.lower() in question_text.lower() for concept in key_concepts[:10]):
                                        quality_score += 1
                                    else:
                                        quality_issues.append("no key concept")
                                    
                                    # Check 3: Options are well-balanced in length (+1)
                                    option_texts = [opt["text"] for opt in q["options"]]
                                    option_lengths = [len(text) for text in option_texts]
                                    length_variance = max(option_lengths) - min(option_lengths)
                                    if length_variance < 50:
                                        quality_score += 1
                                    else:
                                        quality_issues.append(f"unbalanced lengths (±{length_variance})")
                                    
                                    # Check 4: Question is specific, not overly generic (+1)
                                    overly_generic = ["based on the content", "according to the text", "from the material", "mentioned above"]
                                    if not any(phrase in question_text.lower() for phrase in overly_generic):
                                        quality_score += 1
                                    else:
                                        quality_issues.append("generic phrasing")
                                    
                                    # Check 5: No "None/All of the above" in options (bonus check)
                                    bad_options = ["none of the above", "all of the above", "both a and b", "a and b"]
                                    has_bad_option = any(
                                        any(bad in opt["text"].lower() for bad in bad_options)
                                        for opt in q["options"]
                                    )
                                    
                                    # Check 6: Options are meaningful (not too generic)
                                    generic_option_phrases = ["is not mentioned", "not discussed", "unrelated", "none"]
                                    meaningful_options = sum(
                                        1 for opt in option_texts 
                                        if not any(phrase in opt.lower() for phrase in generic_option_phrases)
                                    )
                                    
                                    if not has_bad_option and meaningful_options >= 3:
                                        quality_score += 1
                                        
                                    # Only accept high-quality questions (score >= 3)
                                    if quality_score >= 3:
                                        valid_questions.append(q)
                                        quality_scores.append(quality_score)
                                        
                                        score_emoji = "⭐" * quality_score
                                        print(f"  ✓ Question {idx}: {score_emoji} [{quality_score}/5] '{question_text[:60]}...'")
                                    else:
                                        print(f"  ⚠ Question {idx}: Low quality [{quality_score}/5] - {', '.join(quality_issues)}")
                                else:
                                    print(f"  ✗ Question {idx}: Invalid or too short options")
                            else:
                                print(f"  ✗ Question {idx}: Question text too short")
                        else:
                            print(f"  ✗ Question {idx}: Has {correct_count} correct answers (need exactly 1)")
                    else:
                        print(f"  ✗ Question {idx}: Invalid structure")
                
                # Log quality metrics
                if quality_scores:
                    avg_quality = sum(quality_scores) / len(quality_scores)
                    print(f"📊 Average question quality: {avg_quality:.1f}/5.0")
                    
                    # Quality rating
                    if avg_quality >= 4.5:
                        print(f"🏆 EXCELLENT quality - Exam-ready questions!")
                    elif avg_quality >= 4.0:
                        print(f"✨ HIGH quality - Professional questions")
                    elif avg_quality >= 3.5:
                        print(f"✓ GOOD quality - Acceptable questions")
                    else:
                        print(f"⚠️ MODERATE quality - Consider regenerating")
                
                print(f"✅ {len(valid_questions)} valid questions after validation")
                
                # If AI generated fewer valid questions, use intelligent fallback
                if len(valid_questions) < num_questions:
                    print(f"⚠️ Need {num_questions - len(valid_questions)} more questions")
                    remaining = num_questions - len(valid_questions)
                    additional = generate_enhanced_mcqs(structured_content, remaining, difficulty)
                    valid_questions.extend(additional)
                
                return valid_questions[:num_questions]
            else:
                print(f"❌ No JSON found in AI response")
                print(f"Response preview: {response_text[:200]}")
    except Exception as e:
        print(f"❌ AI generation error: {e}")
        import traceback
        traceback.print_exc()
    
    # Fallback: Generate enhanced questions from structured content
    print(f"📝 Using enhanced fallback method to generate questions")
    return generate_enhanced_mcqs(structured_content, num_questions, difficulty)

def generate_enhanced_mcqs(structured_content: dict, num_questions: int, difficulty: str) -> List[dict]:
    """Enhanced fallback method to generate high-quality MCQs from structured content"""
    
    questions = []
    headings = structured_content.get("headings", [])
    important_paras = structured_content.get("important_paragraphs", [])
    # Ensure key_terms is a list
    key_terms_raw = structured_content.get("key_terms", [])
    key_terms = list(key_terms_raw) if isinstance(key_terms_raw, set) else key_terms_raw
    
    # Question templates based on difficulty
    templates = {
        "easy": [
            ("What is {concept}?", "definition"),
            ("Which of the following describes {concept}?", "description"),
            ("According to the content, what does {concept} refer to?", "reference"),
            ("{concept} is best defined as:", "definition"),
        ],
        "medium": [
            ("How does {concept} relate to the main topic?", "relationship"),
            ("What is the primary purpose of {concept}?", "purpose"),
            ("Which statement best explains {concept}?", "explanation"),
            ("What can be inferred about {concept} from the content?", "inference"),
        ],
        "hard": [
            ("Based on the content, what is the significance of {concept}?", "significance"),
            ("Analyze the role of {concept} in the context discussed:", "analysis"),
            ("What would be the consequence if {concept} were changed?", "evaluation"),
            ("Compare {concept} with related concepts in the content:", "synthesis"),
        ]
    }
    
    current_templates = templates.get(difficulty.lower(), templates["medium"])
    
    # Strategy 1: Create questions from headings (most important)
    for i, heading in enumerate(headings[:num_questions]):
        if len(questions) >= num_questions:
            break
        
        template, q_type = current_templates[i % len(current_templates)]
        question_text = template.replace("{concept}", heading)
        
        # Create contextual answer based on next paragraph
        correct_answer = heading
        if i < len(important_paras):
            # Use content from paragraph as correct answer
            para_words = important_paras[i].split()[:15]
            correct_answer = ' '.join(para_words)
        
        # Generate plausible wrong answers
        wrong_answers = generate_plausible_distractors(heading, important_paras, i, correct_answer)
        
        questions.append({
            "question": question_text,
            "options": [
                {"text": correct_answer, "is_correct": True},
                {"text": wrong_answers[0], "is_correct": False},
                {"text": wrong_answers[1], "is_correct": False},
                {"text": wrong_answers[2], "is_correct": False}
            ],
            "explanation": f"The content discusses {heading.lower()}, which is a key concept in understanding the material."
        })
    
    # Strategy 2: Create questions from key terms
    for i, term in enumerate(key_terms[:num_questions - len(questions)]):
        if len(questions) >= num_questions:
            break
        
        template, q_type = current_templates[i % len(current_templates)]
        question_text = template.replace("{concept}", term)
        
        # Find context for this term in paragraphs
        correct_answer = f"A fundamental concept related to {term}"
        for para in important_paras:
            if term.lower() in para.lower():
                # Extract sentence containing the term
                sentences = para.split('. ')
                for sent in sentences:
                    if term.lower() in sent.lower():
                        correct_answer = sent.strip()
                        break
                break
        
        wrong_answers = generate_plausible_distractors(term, important_paras, len(questions), correct_answer)
        
        questions.append({
            "question": question_text,
            "options": [
                {"text": correct_answer, "is_correct": True},
                {"text": wrong_answers[0], "is_correct": False},
                {"text": wrong_answers[1], "is_correct": False},
                {"text": wrong_answers[2], "is_correct": False}
            ],
            "explanation": f"This relates to {term}, an important term in the content."
        })
    
    # Strategy 3: Create questions from important paragraphs
    para_start = len(questions)
    for i, para in enumerate(important_paras):
        if len(questions) >= num_questions:
            break
        
        # Skip if already used this paragraph
        if i < para_start:
            continue
        
        # Extract main idea from paragraph
        sentences = [s.strip() for s in para.split('. ') if len(s.strip()) > 20]
        if not sentences:
            continue
        
        main_sentence = sentences[0]  # First sentence often contains main idea
        
        # Extract key concept from sentence
        words = main_sentence.split()
        key_words = [w for w in words if len(w) > 4 and w[0].isupper()]
        concept = key_words[0] if key_words else "this concept"
        
        template_idx = (len(questions) + i) % len(current_templates)
        template, q_type = current_templates[template_idx]
        question_text = template.replace("{concept}", concept)
        
        # Use part of the paragraph as correct answer
        correct_answer = main_sentence
        if len(correct_answer) > 100:
            correct_answer = ' '.join(main_sentence.split()[:15]) + "..."
        
        wrong_answers = generate_plausible_distractors(concept, important_paras, i, correct_answer)
        
        questions.append({
            "question": question_text,
            "options": [
                {"text": correct_answer, "is_correct": True},
                {"text": wrong_answers[0], "is_correct": False},
                {"text": wrong_answers[1], "is_correct": False},
                {"text": wrong_answers[2], "is_correct": False}
            ],
            "explanation": f"According to the content: {main_sentence[:150]}..."
        })
    
    # If still need more questions, create general comprehension questions
    while len(questions) < num_questions:
        idx = len(questions)
        questions.append({
            "question": f"Based on the content, which of the following is a key concept discussed?",
            "options": [
                {"text": headings[idx % len(headings)] if headings else "Main concept from the content", "is_correct": True},
                {"text": "This concept is not mentioned in the material", "is_correct": False},
                {"text": "An unrelated topic", "is_correct": False},
                {"text": "A contradictory statement", "is_correct": False}
            ],
            "explanation": "This question tests understanding of the main concepts presented in the content."
        })
    
    return questions[:num_questions]

def generate_plausible_distractors(concept: str, paragraphs: List[str], exclude_idx: int, correct_answer: str) -> List[str]:
    """Generate meaningful, realistic but incorrect answer options from document content"""
    distractors = []
    
    # Strategy 1: Extract related concepts from other paragraphs
    for i, para in enumerate(paragraphs[:10]):
        if i != exclude_idx and len(distractors) < 3:
            # Find sentences or phrases that are related but different
            sentences = [s.strip() for s in para.split('.') if len(s.strip()) > 30]
            for sent in sentences:
                if sent != correct_answer and len(sent) > 25:
                    # Make it similar length to correct answer
                    words = sent.split()
                    target_length = len(correct_answer.split())
                    
                    if len(words) >= target_length - 3:
                        distractor_text = ' '.join(words[:target_length + 2])
                        if distractor_text not in distractors and distractor_text != correct_answer:
                            distractors.append(distractor_text)
                            if len(distractors) >= 3:
                                break
            if len(distractors) >= 3:
                break
    
    # Strategy 2: Create variations based on the concept
    if len(distractors) < 3:
        # Use similar structure but different meaning
        correct_words = correct_answer.split()
        if len(correct_words) > 5:
            # Partial truth - use first part differently
            variant = ' '.join(correct_words[:len(correct_words)//2]) + " in a different context"
            if variant not in distractors:
                distractors.append(variant)
    
    # Strategy 3: Use content from document but for different concepts
    if len(distractors) < 3:
        for para in paragraphs[exclude_idx+1:exclude_idx+5]:
            if len(distractors) >= 3:
                break
            words = para.split()[:len(correct_answer.split()) + 1]
            if len(words) > 5:
                alt_text = ' '.join(words)
                if alt_text != correct_answer and alt_text not in distractors:
                    distractors.append(alt_text)
    
    # Ensure we have exactly 3 distractors with similar length
    while len(distractors) < 3:
        # Last resort: create meaningful but incorrect statement
        distractors.append(f"Related to {concept} but represents a different aspect of the topic")
    
    return distractors[:3]

# Routes
@app.get("/")
async def root():
    return {"message": "Quizzify API is running"}

@app.post("/register")
async def register(user: UserRegister):
    # Check if email already exists in users collection
    existing_email = await db.users.find_one({"email": user.email})
    if existing_email:
        raise HTTPException(status_code=400, detail="Email already registered")
    
    # Check if username already exists
    existing_username = await db.users.find_one({"username": user.username})
    if existing_username:
        raise HTTPException(status_code=400, detail="Username already taken")
    
    try:
        # Generate OTP
        otp = generate_otp()
        otp_expiry = datetime.utcnow() + timedelta(minutes=10)
        
        # Hash password
        hashed_pwd = get_password_hash(user.password)
        
        # Store pending registration with OTP
        pending_user = {
            "username": user.username,
            "email": user.email,
            "full_name": user.full_name,
            "hashed_password": hashed_pwd,
            "otp": otp,
            "otp_expiry": otp_expiry,
            "created_at": datetime.utcnow(),
            "verified": False
        }
        
        # Delete any existing pending registration for this email
        await db.pending_users.delete_many({"email": user.email})
        
        # Insert pending user
        await db.pending_users.insert_one(pending_user)
        
        # Send OTP email
        email_sent = await send_otp_email(user.email, otp)
        
        print(f"📧 OTP generated for {user.email}: {otp} (expires in 10 minutes)")
        
        return {
            "message": "OTP sent to your email. Please verify to complete registration.",
            "email": user.email,
            "email_sent": email_sent
        }
    except Exception as e:
        print(f"Registration error: {e}")
        raise HTTPException(status_code=500, detail="Registration failed. Please try again.")

@app.post("/verify-otp")
async def verify_otp(verification: OTPVerification):
    """Verify OTP and complete user registration"""
    try:
        # Find pending user
        pending_user = await db.pending_users.find_one({"email": verification.email})
        
        if not pending_user:
            raise HTTPException(status_code=404, detail="No pending registration found for this email")
        
        # Check if OTP expired
        if datetime.utcnow() > pending_user["otp_expiry"]:
            await db.pending_users.delete_one({"email": verification.email})
            raise HTTPException(status_code=400, detail="OTP has expired. Please register again.")
        
        # Verify OTP
        if pending_user["otp"] != verification.otp:
            raise HTTPException(status_code=400, detail="Invalid OTP. Please try again.")
        
        # OTP is valid - create the user
        user_dict = {
            "username": pending_user["username"],
            "email": pending_user["email"],
            "full_name": pending_user.get("full_name"),
            "hashed_password": pending_user["hashed_password"],
            "created_at": datetime.utcnow(),
            "verified": True
        }
        
        result = await db.users.insert_one(user_dict)
        
        # Delete pending user
        await db.pending_users.delete_one({"email": verification.email})
        
        print(f"✅ Email verified and user created: {verification.email}")
        
        # Create access token for auto-login
        access_token = create_access_token(data={"sub": verification.email})
        
        return {
            "message": "Email verified successfully! You can now login.",
            "access_token": access_token,
            "token_type": "bearer"
        }
        
    except HTTPException:
        raise
    except Exception as e:
        print(f"OTP verification error: {e}")
        raise HTTPException(status_code=500, detail="Verification failed. Please try again.")

@app.post("/resend-otp")
async def resend_otp(data: ResendOTP):
    """Resend OTP to email"""
    try:
        # Find pending user
        pending_user = await db.pending_users.find_one({"email": data.email})
        
        if not pending_user:
            raise HTTPException(status_code=404, detail="No pending registration found for this email")
        
        # Generate new OTP
        otp = generate_otp()
        otp_expiry = datetime.utcnow() + timedelta(minutes=10)
        
        # Update pending user with new OTP
        await db.pending_users.update_one(
            {"email": data.email},
            {"$set": {"otp": otp, "otp_expiry": otp_expiry}}
        )
        
        # Send OTP email
        email_sent = await send_otp_email(data.email, otp)
        
        print(f"📧 New OTP generated for {data.email}: {otp}")
        
        return {
            "message": "New OTP sent to your email",
            "email_sent": email_sent
        }
        
    except HTTPException:
        raise
    except Exception as e:
        print(f"Resend OTP error: {e}")
        raise HTTPException(status_code=500, detail="Failed to resend OTP. Please try again.")

@app.post("/token", response_model=Token)
async def login(form_data: OAuth2PasswordRequestForm = Depends()):
    user = await db.users.find_one({"email": form_data.username})
    if not user or not verify_password(form_data.password, user["hashed_password"]):
        raise HTTPException(
            status_code=status.HTTP_401_UNAUTHORIZED,
            detail="Incorrect email or password",
            headers={"WWW-Authenticate": "Bearer"},
        )
    
    access_token = create_access_token(data={"sub": user["email"]})
    return {"access_token": access_token, "token_type": "bearer"}

@app.post("/login", response_model=Token)
async def login_json(user_data: UserLogin):
    user = await db.users.find_one({"email": user_data.email})
    
    if not user:
        print(f"Login failed: User not found with email {user_data.email}")
        raise HTTPException(
            status_code=status.HTTP_401_UNAUTHORIZED,
            detail="Incorrect email or password"
        )
    
    password_valid = verify_password(user_data.password, user["hashed_password"])
    print(f"Login attempt for {user_data.email}: Password valid = {password_valid}")
    
    if not password_valid:
        raise HTTPException(
            status_code=status.HTTP_401_UNAUTHORIZED,
            detail="Incorrect email or password"
        )
    
    access_token = create_access_token(data={"sub": user["email"]})
    return {"access_token": access_token, "token_type": "bearer"}

@app.get("/me")
async def get_me(current_user: dict = Depends(get_current_user)):
    return {
        "email": current_user["email"],
        "username": current_user["username"],
        "full_name": current_user.get("full_name")
    }

@app.put("/me")
async def update_profile(
    profile_data: UpdateProfile,
    current_user: dict = Depends(get_current_user)
):
    """Update user profile"""
    # Check if username is taken by another user
    existing_user = await db.users.find_one({
        "username": profile_data.username,
        "email": {"$ne": current_user["email"]}
    })
    if existing_user:
        raise HTTPException(status_code=400, detail="Username already taken")
    
    # Update user
    await db.users.update_one(
        {"email": current_user["email"]},
        {"$set": {
            "username": profile_data.username,
            "full_name": profile_data.full_name
        }}
    )
    
    # Get updated user
    updated_user = await db.users.find_one({"email": current_user["email"]})
    
    return {
        "message": "Profile updated successfully",
        "user": {
            "email": updated_user["email"],
            "username": updated_user["username"],
            "full_name": updated_user.get("full_name")
        }
    }

@app.put("/me/password")
async def update_password(
    password_data: UpdatePassword,
    current_user: dict = Depends(get_current_user)
):
    """Update user password"""
    # Verify current password
    if not verify_password(password_data.current_password, current_user["hashed_password"]):
        raise HTTPException(
            status_code=status.HTTP_401_UNAUTHORIZED,
            detail="Current password is incorrect"
        )
    
    # Update password
    new_hashed_password = get_password_hash(password_data.new_password)
    await db.users.update_one(
        {"email": current_user["email"]},
        {"$set": {"hashed_password": new_hashed_password}}
    )
    
    return {"message": "Password updated successfully"}

@app.post("/upload-and-generate")
async def upload_and_generate(
    file: UploadFile = File(...),
    num_questions: int = 5,
    difficulty: str = "medium",
    quiz_name: str = "",
    current_user: dict = Depends(get_current_user)
):
    """Upload a file and generate high-quality MCQs focused on important concepts"""
    
    # Read file content
    content = await file.read()
    
    # Extract text based on file type
    filename = file.filename.lower()
    
    try:
        print(f"\n📄 Processing file: {file.filename}")
        print(f"🎯 Target: {num_questions} questions at {difficulty} difficulty")
        
        # Extract structured content based on file type
        if filename.endswith('.pdf'):
            structured_content = extract_text_from_pdf(content)
            print(f"📊 Extracted from PDF: {len(structured_content.get('headings', []))} headings, {len(structured_content.get('important_paragraphs', []))} important paragraphs")
        elif filename.endswith('.pptx') or filename.endswith('.ppt'):
            structured_content = extract_text_from_pptx(content)
            print(f"📊 Extracted from PPTX: {len(structured_content.get('headings', []))} slide titles, {len(structured_content.get('important_paragraphs', []))} content blocks")
        elif filename.endswith('.docx') or filename.endswith('.doc'):
            structured_content = extract_text_from_docx(content)
            print(f"📊 Extracted from DOCX: {len(structured_content.get('headings', []))} headings, {len(structured_content.get('key_terms', []))} key terms")
        elif filename.endswith('.txt'):
            # For plain text, create basic structure
            text = content.decode('utf-8')
            structured_content = {
                "full_text": text,
                "headings": [],
                "important_paragraphs": [p.strip() for p in text.split('\n\n') if len(p.strip()) > 50][:20],
                "key_terms": []
            }
            print(f"📊 Extracted from TXT: {len(structured_content['important_paragraphs'])} paragraphs")
        else:
            raise HTTPException(status_code=400, detail="Unsupported file type. Please upload PDF, DOCX, PPTX, or TXT files.")
        
        # Validate extracted content
        full_text = structured_content.get("full_text", "")
        if not full_text.strip():
            raise HTTPException(status_code=400, detail="No text could be extracted from the file. Please ensure the file contains readable text.")
        
        if len(full_text.strip()) < 100:
            raise HTTPException(status_code=400, detail="The extracted content is too short. Please provide a document with more substantial content.")
        
        print(f"✅ Content extracted successfully ({len(full_text)} characters)")
        
        # Generate high-quality MCQs using enhanced AI
        questions = await generate_mcqs_with_ai(structured_content, num_questions, difficulty)
        
        if not questions or len(questions) == 0:
            raise HTTPException(status_code=500, detail="Failed to generate questions from the content. Please try again or use a different document.")
        
        print(f"✅ Generated {len(questions)} high-quality questions")
        
        # Use custom quiz name if provided, otherwise use filename
        title = quiz_name.strip() if quiz_name.strip() else f"Quiz from {file.filename}"
        
        # Save quiz to database with metadata
        quiz_data = {
            "title": title,
            "questions": questions,
            "created_by": current_user["email"],
            "created_at": datetime.utcnow(),
            "source_file": file.filename,
            "difficulty": difficulty,
            "num_concepts": len(structured_content.get("headings", [])),
            "content_length": len(full_text)
        }
        
        result = await db.quizzes.insert_one(quiz_data)
        quiz_data["id"] = str(result.inserted_id)
        del quiz_data["_id"]  # Remove the ObjectId before returning
        
        return {
            "message": "High-quality quiz generated successfully from important concepts",
            "quiz": quiz_data,
            "num_questions": len(questions),
            "insights": {
                "headings_found": len(structured_content.get("headings", [])),
                "key_concepts": len(structured_content.get("key_terms", [])),
                "content_analyzed": f"{len(full_text)} characters"
            }
        }
        
    except HTTPException:
        raise
    except Exception as e:
        import traceback
        print(f"❌ Error processing file: {str(e)}")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Error processing file: {str(e)}")

@app.get("/quizzes")
async def get_quizzes(current_user: dict = Depends(get_current_user)):
    """Get all quizzes created by the user"""
    quizzes = await db.quizzes.find({"created_by": current_user["email"]}).to_list(100)
    
    for quiz in quizzes:
        quiz["id"] = str(quiz["_id"])
        del quiz["_id"]
    
    return {"quizzes": quizzes}

@app.get("/quizzes/all/public")
async def get_all_public_quizzes(current_user: dict = Depends(get_current_user)):
    """Get all quizzes from all users (public quizzes)"""
    quizzes = await db.quizzes.find({}).sort("created_at", -1).to_list(100)
    
    for quiz in quizzes:
        quiz["id"] = str(quiz["_id"])
        del quiz["_id"]
    
    return {"quizzes": quizzes}

@app.get("/quizzes/{quiz_id}")
async def get_quiz(quiz_id: str, current_user: dict = Depends(get_current_user)):
    """Get a specific quiz"""
    from bson import ObjectId
    
    try:
        quiz = await db.quizzes.find_one({"_id": ObjectId(quiz_id)})
        if not quiz:
            raise HTTPException(status_code=404, detail="Quiz not found")
        
        quiz["id"] = str(quiz["_id"])
        del quiz["_id"]
        
        return quiz
    except Exception as e:
        raise HTTPException(status_code=400, detail="Invalid quiz ID")

@app.delete("/quizzes/{quiz_id}")
async def delete_quiz(quiz_id: str, current_user: dict = Depends(get_current_user)):
    """Delete a quiz"""
    from bson import ObjectId
    
    try:
        result = await db.quizzes.delete_one({
            "_id": ObjectId(quiz_id),
            "created_by": current_user["email"]
        })
        
        if result.deleted_count == 0:
            raise HTTPException(status_code=404, detail="Quiz not found")
        
        return {"message": "Quiz deleted successfully"}
    except Exception as e:
        raise HTTPException(status_code=400, detail="Invalid quiz ID")

@app.put("/quizzes/{quiz_id}")
async def update_quiz(quiz_id: str, quiz_data: dict, current_user: dict = Depends(get_current_user)):
    """Update an existing quiz"""
    from bson import ObjectId
    
    try:
        # Verify quiz exists and user owns it
        existing_quiz = await db.quizzes.find_one({
            "_id": ObjectId(quiz_id),
            "created_by": current_user["email"]
        })
        
        if not existing_quiz:
            raise HTTPException(status_code=404, detail="Quiz not found or you don't have permission to edit it")
        
        # Update quiz
        update_data = {
            "title": quiz_data.get("title", existing_quiz["title"]),
            "questions": quiz_data.get("questions", existing_quiz["questions"]),
            "updated_at": datetime.utcnow()
        }
        
        await db.quizzes.update_one(
            {"_id": ObjectId(quiz_id)},
            {"$set": update_data}
        )
        
        return {"message": "Quiz updated successfully", "quiz_id": quiz_id}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Failed to update quiz: {str(e)}")

class ManualQuizRequest(BaseModel):
    title: str
    questions: List[MCQuestion]

@app.post("/quizzes/create-manual")
async def create_manual_quiz(
    quiz_request: ManualQuizRequest,
    current_user: dict = Depends(get_current_user)
):
    """Create a quiz manually"""
    
    if not quiz_request.title.strip():
        raise HTTPException(status_code=400, detail="Quiz title is required")
    
    if not quiz_request.questions or len(quiz_request.questions) == 0:
        raise HTTPException(status_code=400, detail="At least one question is required")
    
    # Validate questions
    for i, question in enumerate(quiz_request.questions):
        if not question.question.strip():
            raise HTTPException(status_code=400, detail=f"Question {i+1}: Question text is required")
        
        if len(question.options) < 2:
            raise HTTPException(status_code=400, detail=f"Question {i+1}: At least 2 options are required")
        
        correct_count = sum(1 for opt in question.options if opt.is_correct)
        if correct_count != 1:
            raise HTTPException(status_code=400, detail=f"Question {i+1}: Exactly one correct answer is required")
    
    # Save quiz to database
    quiz_data = {
        "title": quiz_request.title,
        "questions": [q.dict() for q in quiz_request.questions],
        "created_by": current_user["email"],
        "created_at": datetime.utcnow(),
        "source": "manual"
    }
    
    result = await db.quizzes.insert_one(quiz_data)
    quiz_data["id"] = str(result.inserted_id)
    del quiz_data["_id"]
    
    return {
        "message": "Quiz created successfully",
        "quiz": quiz_data
    }

# ============= ROOM ENDPOINTS =============

def generate_room_code():
    """Generate a unique 6-character room code"""
    import random
    import string
    return ''.join(random.choices(string.ascii_uppercase + string.digits, k=6))

class CreateRoomRequest(BaseModel):
    quiz_id: str
    title: str
    description: Optional[str] = None
    settings: RoomSettings

@app.post("/rooms/create")
async def create_room(
    room_request: CreateRoomRequest,
    current_user: dict = Depends(get_current_user)
):
    """Create a new multiplayer room"""
    from bson import ObjectId
    
    # Verify quiz exists and belongs to user
    try:
        quiz = await db.quizzes.find_one({
            "_id": ObjectId(room_request.quiz_id),
            "created_by": current_user["email"]
        })
        if not quiz:
            raise HTTPException(status_code=404, detail="Quiz not found")
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Invalid quiz ID: {str(e)}")
    
    # Generate unique room code
    room_code = generate_room_code()
    while await db.rooms.find_one({"room_code": room_code}):
        room_code = generate_room_code()
    
    # Create room (host is not included in participants list)
    room_data = {
        "title": room_request.title,
        "description": room_request.description,
        "quiz_id": room_request.quiz_id,
        "host_email": current_user["email"],
        "room_code": room_code,
        "settings": room_request.settings.dict(),
        "status": "waiting",
        "participants": [],  # Host is not a participant
        "max_participants": 50,
        "created_at": datetime.utcnow(),
        "started_at": None
    }
    
    result = await db.rooms.insert_one(room_data)
    room_data["id"] = str(result.inserted_id)
    del room_data["_id"]
    
    return {
        "message": "Room created successfully",
        "room": room_data
    }

@app.get("/rooms")
async def get_rooms(current_user: dict = Depends(get_current_user)):
    """Get all available rooms"""
    rooms = await db.rooms.find({"status": {"$in": ["waiting", "active"]}}).to_list(100)
    
    for room in rooms:
        room["id"] = str(room["_id"])
        del room["_id"]
        
        # Get quiz title
        from bson import ObjectId
        quiz = await db.quizzes.find_one({"_id": ObjectId(room["quiz_id"])})
        if quiz:
            room["quiz_title"] = quiz["title"]
        
        # Get host username
        host = await db.users.find_one({"email": room["host_email"]})
        if host:
            room["host_username"] = host["username"]
    
    return {"rooms": rooms}

@app.get("/rooms/my-rooms")
async def get_my_rooms(current_user: dict = Depends(get_current_user)):
    """Get rooms created by the current user"""
    print(f"🔍 Fetching rooms for user: {current_user['email']}")
    rooms = await db.rooms.find({"host_email": current_user["email"]}).to_list(100)
    print(f"📊 Found {len(rooms)} rooms for {current_user['email']}")
    
    for room in rooms:
        room["id"] = str(room["_id"])
        del room["_id"]
        
        # Get quiz title
        from bson import ObjectId
        quiz = await db.quizzes.find_one({"_id": ObjectId(room["quiz_id"])})
        if quiz:
            room["quiz_title"] = quiz["title"]
    
    return {"rooms": rooms}

@app.get("/rooms/{room_id}")
async def get_room(room_id: str, current_user: dict = Depends(get_current_user)):
    """Get room details"""
    from bson import ObjectId
    
    try:
        room = await db.rooms.find_one({"_id": ObjectId(room_id)})
        if not room:
            raise HTTPException(status_code=404, detail="Room not found")
        
        room["id"] = str(room["_id"])
        del room["_id"]
        
        # Get quiz details
        quiz = await db.quizzes.find_one({"_id": ObjectId(room["quiz_id"])})
        if quiz:
            room["quiz"] = {
                "id": str(quiz["_id"]),
                "title": quiz["title"],
                "num_questions": len(quiz.get("questions", []))
            }
        
        # Get participants details
        participants_details = []
        for participant_email in room.get("participants", []):
            user = await db.users.find_one({"email": participant_email})
            if user:
                participants_details.append({
                    "email": participant_email,
                    "username": user["username"],
                    "full_name": user.get("full_name")
                })
        room["participants_details"] = participants_details
        
        return room
    except Exception as e:
        raise HTTPException(status_code=400, detail="Invalid room ID")

@app.post("/rooms/join")
async def join_room(
    join_request: JoinRoomRequest,
    current_user: dict = Depends(get_current_user)
):
    """
    Join a room using room code (PARTICIPANTS ONLY)
    Host cannot join their own room as participant
    Uses atomic $addToSet to prevent duplicates
    """
    from bson import ObjectId
    
    room = await db.rooms.find_one({"room_code": join_request.room_code.upper()})
    
    if not room:
        raise HTTPException(status_code=404, detail="Room not found. Please check the room code.")
    
    # Prevent host from joining their own room as participant
    if room["host_email"] == current_user["email"]:
        print(f"⚠️  Host {current_user['email']} tried to join their own room {room['room_code']}")
        raise HTTPException(
            status_code=400, 
            detail="You are the host of this room. Hosts cannot join as participants."
        )
    
    # Check room status
    if room["status"] == "completed":
        raise HTTPException(status_code=400, detail="This room has already completed")
    
    # Check room capacity
    if len(room.get("participants", [])) >= room.get("max_participants", 50):
        raise HTTPException(status_code=400, detail="Room is full")
    
    # Atomic add to prevent duplicate participants (race condition safe)
    result = await db.rooms.update_one(
        {
            "_id": ObjectId(room["_id"]),
            "status": {"$in": ["waiting", "active"]},  # Can join waiting or active rooms
            "participants": {"$ne": current_user["email"]}  # Only if not already participant
        },
        {
            "$addToSet": {"participants": current_user["email"]},  # Atomic add (no duplicates)
            "$set": {"updated_at": datetime.utcnow()}
        }
    )
    
    # If no document was modified, user was already a participant or room state changed
    if result.modified_count == 0:
        # Check if already participant
        refreshed_room = await db.rooms.find_one({"_id": ObjectId(room["_id"])})
        if current_user["email"] in refreshed_room.get("participants", []):
            print(f"ℹ️  User {current_user['email']} already in room {room['room_code']}")
            return {
                "message": "You are already in this room",
                "room_id": str(room["_id"]),
                "room_code": room["room_code"],
                "already_joined": True
            }
        else:
            raise HTTPException(
                status_code=400, 
                detail="Could not join room. Room may be full or closed."
            )
    
    print(f"✅ User {current_user['email']} joined room {room['room_code']}")
    return {
        "message": "Joined room successfully",
        "room_id": str(room["_id"]),
        "room_code": room["room_code"],
        "already_joined": False
    }

@app.post("/rooms/{room_id}/start")
async def start_room(room_id: str, current_user: dict = Depends(get_current_user)):
    """
    Start the quiz room (HOST ONLY - STRICT VALIDATION)
    Only the room creator can start the quiz
    """
    from bson import ObjectId
    
    try:
        # Fetch room with explicit host validation
        room = await db.rooms.find_one({"_id": ObjectId(room_id)})
        if not room:
            raise HTTPException(status_code=404, detail="Room not found")
        
        print(f"🎮 Start room request from {current_user['email']} for room {room_id}")
        print(f"   Room host: {room.get('host_email')}")
        print(f"   Current user: {current_user['email']}")
        print(f"   Room status: {room.get('status')}")
        
        # STRICT HOST CHECK - exact email match (case-sensitive, trimmed)
        room_host = str(room["host_email"]).strip()
        current_email = str(current_user["email"]).strip()
        
        if room_host != current_email:
            print(f"❌ NOT HOST: '{current_email}' tried to start room owned by '{room_host}'")
            raise HTTPException(
                status_code=403, 
                detail=f"Only the host can start the room. Host: {room_host}, You: {current_email}"
            )
        
        # Check room status
        if room["status"] != "waiting":
            print(f"⚠️  Room already started (status: {room['status']})")
            raise HTTPException(
                status_code=400, 
                detail=f"Room cannot be started. Current status: {room['status']}"
            )
        
        # Atomic update to prevent race conditions
        result = await db.rooms.update_one(
            {
                "_id": ObjectId(room_id),
                "host_email": current_email,  # Double-check host in query
                "status": "waiting"  # Only update if still waiting
            },
            {
                "$set": {
                    "status": "active",
                    "started_at": datetime.utcnow()
                }
            }
        )
        
        if result.modified_count == 0:
            raise HTTPException(
                status_code=400, 
                detail="Room could not be started. It may have already been started by another request."
            )
        
        print(f"✅ Room {room_id} started successfully by host {current_email}!")
        return {
            "message": "Room started successfully", 
            "status": "active",
            "started_by": current_email
        }
        
    except HTTPException:
        raise
    except Exception as e:
        print(f"❌ Error starting room: {str(e)}")
        raise HTTPException(status_code=400, detail=f"Invalid room ID or server error: {str(e)}")

@app.post("/rooms/{room_id}/complete")
async def complete_room(room_id: str, current_user: dict = Depends(get_current_user)):
    """
    End/Complete the quiz room (HOST ONLY - STRICT VALIDATION)
    Only the room creator can end the quiz
    """
    from bson import ObjectId
    
    try:
        # Fetch room with explicit host validation
        room = await db.rooms.find_one({"_id": ObjectId(room_id)})
        if not room:
            raise HTTPException(status_code=404, detail="Room not found")
        
        # STRICT HOST CHECK - exact email match (case-sensitive, trimmed)
        room_host = str(room["host_email"]).strip()
        current_email = str(current_user["email"]).strip()
        
        print(f"🏁 End room request from {current_email} for room {room_id}")
        print(f"   Room host: {room_host}")
        print(f"   Room status: {room.get('status')}")
        
        if room_host != current_email:
            print(f"❌ NOT HOST: '{current_email}' tried to end room owned by '{room_host}'")
            raise HTTPException(
                status_code=403, 
                detail=f"Only the host can end the room. Host: {room_host}, You: {current_email}"
            )
        
        # Check if room can be ended
        if room["status"] == "completed":
            raise HTTPException(status_code=400, detail="Room is already completed")
        
        # Atomic update to prevent race conditions
        result = await db.rooms.update_one(
            {
                "_id": ObjectId(room_id),
                "host_email": current_email,  # Double-check host in query
                "status": {"$ne": "completed"}  # Only update if not already completed
            },
            {
                "$set": {
                    "status": "completed",
                    "completed_at": datetime.utcnow()
                }
            }
        )
        
        if result.modified_count == 0:
            raise HTTPException(
                status_code=400, 
                detail="Room could not be ended. It may have already been ended."
            )
        
        print(f"✅ Room {room_id} ended successfully by host {current_email}!")
        return {
            "message": "Room ended successfully",
            "status": "completed",
            "ended_by": current_email
        }
        
    except HTTPException:
        raise
    except Exception as e:
        print(f"❌ Error ending room: {str(e)}")
        raise HTTPException(status_code=400, detail=f"Invalid room ID or server error: {str(e)}")

@app.get("/rooms/{room_id}/leaderboard")
async def get_room_leaderboard(room_id: str, current_user: dict = Depends(get_current_user)):
    """Get real-time leaderboard for a room with scores and completion times"""
    from bson import ObjectId
    
    try:
        room = await db.rooms.find_one({"_id": ObjectId(room_id)})
        if not room:
            raise HTTPException(status_code=404, detail="Room not found")
        
        # Check if user is a participant or host
        if current_user["email"] not in room["participants"] and current_user["email"] != room["host_email"]:
            raise HTTPException(status_code=403, detail="You are not a participant in this room")
        
        quiz_id = room.get("quiz_id")
        if not quiz_id:
            return {"leaderboard": [], "total_participants": len(room.get("participants", []))}
        
        # Get all attempts for this quiz from room participants
        leaderboard = []
        for participant_email in room.get("participants", []):
            user = await db.users.find_one({"email": participant_email})
            if not user:
                continue
            
            # Find the best attempt for this user
            attempts = await db.quiz_attempts.find({
                "quiz_id": str(quiz_id),
                "student_email": participant_email
            }).sort("score", -1).limit(1).to_list(length=1)
            
            if attempts:
                attempt = attempts[0]
                # Calculate total time taken (sum of time per question)
                total_time = sum(attempt.get("time_per_question", {}).values())
                
                leaderboard.append({
                    "email": participant_email,
                    "username": user["username"],
                    "full_name": user.get("full_name", ""),
                    "score": attempt.get("score", 0),
                    "max_score": attempt.get("max_score", 0),
                    "percentage": round((attempt.get("score", 0) / attempt.get("max_score", 1)) * 100, 1) if attempt.get("max_score", 0) > 0 else 0,
                    "time_taken": round(total_time, 1),
                    "completed_at": attempt.get("submitted_at"),
                    "correct_answers": attempt.get("correct_answers", 0),
                    "total_questions": len(attempt.get("responses", [])),
                    "has_submitted": True
                })
            else:
                # User hasn't submitted yet
                leaderboard.append({
                    "email": participant_email,
                    "username": user["username"],
                    "full_name": user.get("full_name", ""),
                    "score": 0,
                    "max_score": 0,
                    "percentage": 0,
                    "time_taken": 0,
                    "completed_at": None,
                    "correct_answers": 0,
                    "total_questions": 0,
                    "has_submitted": False
                })
        
        # Sort by score (desc), then by time_taken (asc) - faster time wins in case of tie
        leaderboard.sort(key=lambda x: (-x["score"], x["time_taken"]))
        
        # Add rank
        for idx, entry in enumerate(leaderboard, 1):
            entry["rank"] = idx
        
        return {
            "leaderboard": leaderboard,
            "total_participants": len(room.get("participants", [])),
            "room_status": room.get("status", "waiting")
        }
    except HTTPException:
        raise
    except Exception as e:
        print(f"Error getting leaderboard: {str(e)}")
        raise HTTPException(status_code=400, detail=f"Failed to get leaderboard: {str(e)}")

# Public Quiz Endpoints (No Authentication Required)

@app.get("/public/quiz/{quiz_id}")
async def get_public_quiz(quiz_id: str):
    """Get quiz for public access (no authentication required)"""
    from bson import ObjectId
    
    try:
        quiz = await db.quizzes.find_one({"_id": ObjectId(quiz_id)})
        if not quiz:
            raise HTTPException(status_code=404, detail="Quiz not found")
        
        # Remove sensitive information
        quiz["id"] = str(quiz["_id"])
        del quiz["_id"]
        del quiz["created_by"]  # Don't expose creator info
        
        return {"quiz": quiz}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=400, detail="Invalid quiz ID")

@app.post("/public/quiz/submit")
async def submit_quiz_public(submission: QuizSubmission):
    """Submit quiz answers with detailed scoring (no authentication required)"""
    from bson import ObjectId
    
    try:
        print(f"📥 Received quiz submission:")
        print(f"   Quiz ID: {submission.quiz_id}")
        print(f"   Student: {submission.student_name} ({submission.student_email})")
        print(f"   Answers: {len(submission.answers)} questions answered")
        
        # Get the quiz
        quiz = await db.quizzes.find_one({"_id": ObjectId(submission.quiz_id)})
        if not quiz:
            raise HTTPException(status_code=404, detail="Quiz not found")
        
        # Detailed scoring analysis
        score = 0
        total_questions = len(quiz["questions"])
        correct_answers = []
        incorrect_answers = []
        unanswered = []
        question_details = []
        
        for question_index, question in enumerate(quiz["questions"]):
            user_answers = submission.answers.get(str(question_index), [])
            
            # Find correct answer indices
            correct_indices = []
            for option_index, option in enumerate(question["options"]):
                if option["is_correct"]:
                    correct_indices.append(option_index)
            
            # Determine if answer is correct
            is_correct = set(user_answers) == set(correct_indices) if user_answers else False
            is_answered = len(user_answers) > 0
            
            if not is_answered:
                unanswered.append(question_index)
            elif is_correct:
                correct_answers.append(question_index)
                score += 1
            else:
                incorrect_answers.append(question_index)
            
            # Store question details for review
            question_details.append({
                "question_index": question_index,
                "question": question["question"],
                "options": [opt["text"] for opt in question["options"]],
                "correct_indices": correct_indices,
                "user_answers": user_answers,
                "is_correct": is_correct,
                "is_answered": is_answered,
                "explanation": question.get("explanation", "")
            })
        
        percentage = round((score / total_questions) * 100, 1)
        
        # Save detailed attempt to database
        attempt_data = {
            "quiz_id": submission.quiz_id,
            "quiz_title": quiz["title"],
            "student_name": submission.student_name,
            "student_email": submission.student_email,
            "answers": submission.answers,
            "score": score,
            "total_questions": total_questions,
            "max_score": total_questions,  # Add max_score field for analytics
            "percentage": percentage,
            "time_taken": submission.time_taken,
            "time_per_question": submission.time_per_question,
            "correct_answers": correct_answers,
            "incorrect_answers": incorrect_answers,
            "unanswered": unanswered,
            "question_details": question_details,
            "submitted_at": datetime.utcnow()
        }
        
        result = await db.quiz_attempts.insert_one(attempt_data)
        attempt_id = str(result.inserted_id)
        
        print(f"✅ Quiz attempt saved successfully!")
        print(f"   Attempt ID: {attempt_id}")
        print(f"   Score: {score}/{total_questions} ({percentage}%)")
        
        return {
            "message": "Quiz submitted successfully",
            "attempt_id": attempt_id,
            "score": score,
            "total_questions": total_questions,
            "percentage": percentage,
            "time_taken": submission.time_taken,
            "correct_count": len(correct_answers),
            "incorrect_count": len(incorrect_answers),
            "unanswered_count": len(unanswered),
            "time_formatted": format_time(submission.time_taken)
        }
        
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error submitting quiz: {str(e)}")

@app.get("/public/attempt/{attempt_id}")
async def get_attempt_details(attempt_id: str):
    """Get detailed attempt results for answer review (no authentication required)"""
    from bson import ObjectId
    
    try:
        attempt = await db.quiz_attempts.find_one({"_id": ObjectId(attempt_id)})
        if not attempt:
            raise HTTPException(status_code=404, detail="Attempt not found")
        
        # Format the response
        attempt["id"] = str(attempt["_id"])
        del attempt["_id"]
        
        return {
            "attempt": attempt
        }
        
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=400, detail="Invalid attempt ID")

@app.get("/my-attempts")
async def get_my_attempts(current_user: dict = Depends(get_current_user)):
    """Get all quiz attempts by the current user"""
    try:
        # Find all attempts by this user's email
        attempts_cursor = db.quiz_attempts.find({
            "student_email": current_user["email"]
        }).sort("submitted_at", -1)
        
        attempts = await attempts_cursor.to_list(length=100)
        
        # Format attempts with quiz details
        formatted_attempts = []
        for attempt in attempts:
            # Get quiz title
            from bson import ObjectId
            quiz = await db.quizzes.find_one({"_id": ObjectId(attempt["quiz_id"])})
            quiz_title = quiz.get("title", "Unknown Quiz") if quiz else "Unknown Quiz"
            
            formatted_attempts.append({
                "id": str(attempt["_id"]),
                "quiz_id": attempt["quiz_id"],
                "quiz_title": quiz_title,
                "score": attempt.get("score", 0),
                "total_questions": attempt.get("total_questions", 0),
                "max_score": attempt.get("total_questions", 0),
                "percentage": attempt.get("percentage", 0),
                "time_taken": attempt.get("time_taken", 0),
                "completed_at": attempt.get("submitted_at"),
                "answers": attempt.get("answers", {}),
                "question_details": attempt.get("question_details", [])
            })
        
        return {"attempts": formatted_attempts}
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error fetching attempts: {str(e)}")

@app.delete("/rooms/{room_id}")
async def delete_room(room_id: str, current_user: dict = Depends(get_current_user)):
    """Delete a room (host only)"""
    from bson import ObjectId
    
    try:
        result = await db.rooms.delete_one({
            "_id": ObjectId(room_id),
            "host_email": current_user["email"]
        })
        
        if result.deleted_count == 0:
            raise HTTPException(status_code=404, detail="Room not found")
        
        return {"message": "Room deleted successfully"}
    except Exception as e:
        raise HTTPException(status_code=400, detail="Invalid room ID")

# ==================== Export/Share Endpoints ====================

@app.get("/quizzes/{quiz_id}/attempts")
async def get_quiz_attempts(quiz_id: str, current_user: dict = Depends(get_current_user)):
    """Get all attempts for a quiz (teacher only)"""
    from bson import ObjectId
    
    try:
        # Verify quiz ownership
        quiz = await db.quizzes.find_one({
            "_id": ObjectId(quiz_id),
            "created_by": current_user["email"]
        })
        
        if not quiz:
            raise HTTPException(status_code=404, detail="Quiz not found or unauthorized")
        
        # Get all attempts for this quiz
        attempts_cursor = db.quiz_attempts.find({"quiz_id": quiz_id}).sort("submitted_at", -1)
        attempts = await attempts_cursor.to_list(length=None)
        
        # Format attempts
        formatted_attempts = []
        for attempt in attempts:
            formatted_attempts.append({
                "id": str(attempt["_id"]),
                "student_name": attempt.get("student_name", "Unknown"),
                "student_email": attempt.get("student_email", "N/A"),
                "score": attempt.get("score", 0),
                "total_questions": attempt.get("total_questions", 0),
                "percentage": attempt.get("percentage", 0),
                "time_taken": attempt.get("time_taken", 0),
                "submitted_at": attempt.get("submitted_at").isoformat() if attempt.get("submitted_at") else None
            })
        
        return {
            "quiz_id": quiz_id,
            "quiz_title": quiz.get("title", "Untitled Quiz"),
            "total_attempts": len(formatted_attempts),
            "attempts": formatted_attempts
        }
        
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Error fetching attempts: {str(e)}")

@app.get("/quizzes/{quiz_id}/export/csv")
async def export_quiz_csv(quiz_id: str, current_user: dict = Depends(get_current_user)):
    """Export quiz attempts as CSV"""
    from bson import ObjectId
    
    try:
        # Verify quiz ownership
        quiz = await db.quizzes.find_one({
            "_id": ObjectId(quiz_id),
            "created_by": current_user["email"]
        })
        
        if not quiz:
            raise HTTPException(status_code=404, detail="Quiz not found or unauthorized")
        
        # Get all attempts
        attempts_cursor = db.quiz_attempts.find({"quiz_id": quiz_id}).sort("submitted_at", -1)
        attempts = await attempts_cursor.to_list(length=None)
        
        # Create CSV in memory
        output = io.StringIO()
        writer = csv.writer(output)
        
        # Write header
        writer.writerow([
            "Student Name",
            "Email",
            "Score",
            "Total Questions",
            "Percentage",
            "Time Taken (seconds)",
            "Time Taken (formatted)",
            "Submitted At"
        ])
        
        # Write data
        for attempt in attempts:
            time_taken = attempt.get("time_taken", 0)
            minutes = time_taken // 60
            seconds = time_taken % 60
            time_formatted = f"{minutes}m {seconds}s"
            
            writer.writerow([
                attempt.get("student_name", "Unknown"),
                attempt.get("student_email", "N/A"),
                attempt.get("score", 0),
                attempt.get("total_questions", 0),
                f"{attempt.get('percentage', 0):.1f}%",
                time_taken,
                time_formatted,
                attempt.get("submitted_at").strftime("%Y-%m-%d %H:%M:%S") if attempt.get("submitted_at") else "N/A"
            ])
        
        # Prepare response
        output.seek(0)
        filename = f"{quiz.get('title', 'quiz').replace(' ', '_')}_attempts.csv"
        
        return StreamingResponse(
            iter([output.getvalue()]),
            media_type="text/csv",
            headers={
                "Content-Disposition": f'attachment; filename="{filename}"'
            }
        )
        
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Error exporting CSV: {str(e)}")

@app.get("/quizzes/{quiz_id}/export/pdf")
async def export_quiz_pdf(quiz_id: str, current_user: dict = Depends(get_current_user)):
    """Export quiz with answers as PDF (for printing)"""
    from bson import ObjectId
    
    try:
        # Verify quiz ownership
        quiz = await db.quizzes.find_one({
            "_id": ObjectId(quiz_id),
            "created_by": current_user["email"]
        })
        
        if not quiz:
            raise HTTPException(status_code=404, detail="Quiz not found or unauthorized")
        
        # Create PDF in memory
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=letter, topMargin=0.75*inch, bottomMargin=0.75*inch)
        story = []
        styles = getSampleStyleSheet()
        
        # Custom styles
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=24,
            textColor=colors.HexColor('#4F46E5'),
            spaceAfter=30,
            alignment=TA_CENTER
        )
        
        question_style = ParagraphStyle(
            'Question',
            parent=styles['Heading2'],
            fontSize=14,
            textColor=colors.HexColor('#1F2937'),
            spaceAfter=12,
            spaceBefore=20
        )
        
        option_style = ParagraphStyle(
            'Option',
            parent=styles['Normal'],
            fontSize=11,
            leftIndent=20,
            spaceAfter=6
        )
        
        correct_style = ParagraphStyle(
            'Correct',
            parent=styles['Normal'],
            fontSize=11,
            leftIndent=20,
            textColor=colors.HexColor('#10B981'),
            spaceAfter=6
        )
        
        explanation_style = ParagraphStyle(
            'Explanation',
            parent=styles['Normal'],
            fontSize=10,
            leftIndent=20,
            textColor=colors.HexColor('#6B7280'),
            spaceAfter=12,
            fontName='Helvetica-Oblique'
        )
        
        # Add title
        story.append(Paragraph(quiz.get("title", "Quiz"), title_style))
        story.append(Paragraph(f"Created: {quiz.get('created_at', datetime.now()).strftime('%B %d, %Y')}", styles['Normal']))
        story.append(Paragraph(f"Total Questions: {len(quiz.get('questions', []))}", styles['Normal']))
        story.append(Spacer(1, 0.3*inch))
        
        # Add questions
        for idx, question in enumerate(quiz.get("questions", []), 1):
            # Question text
            story.append(Paragraph(f"<b>Question {idx}:</b> {question.get('question', '')}", question_style))
            
            # Options
            for opt_idx, option in enumerate(question.get("options", []), 1):
                option_text = option.get("text", "")
                is_correct = option.get("is_correct", False)
                
                if is_correct:
                    story.append(Paragraph(f"<b>{chr(64+opt_idx)}. {option_text} ✓</b>", correct_style))
                else:
                    story.append(Paragraph(f"{chr(64+opt_idx)}. {option_text}", option_style))
            
            # Explanation
            if question.get("explanation"):
                story.append(Spacer(1, 0.1*inch))
                story.append(Paragraph(f"<b>Explanation:</b> {question.get('explanation')}", explanation_style))
            
            # Add some space between questions
            story.append(Spacer(1, 0.2*inch))
        
        # Build PDF
        doc.build(story)
        buffer.seek(0)
        
        filename = f"{quiz.get('title', 'quiz').replace(' ', '_')}_with_answers.pdf"
        
        return StreamingResponse(
            buffer,
            media_type="application/pdf",
            headers={
                "Content-Disposition": f'attachment; filename="{filename}"'
            }
        )
        
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Error exporting PDF: {str(e)}")

@app.put("/quizzes/{quiz_id}/share-settings")
async def update_share_settings(
    quiz_id: str,
    share_settings: QuizShareSettings,
    current_user: dict = Depends(get_current_user)
):
    """Update quiz sharing settings"""
    from bson import ObjectId
    
    try:
        # Verify quiz ownership
        quiz = await db.quizzes.find_one({
            "_id": ObjectId(quiz_id),
            "created_by": current_user["email"]
        })
        
        if not quiz:
            raise HTTPException(status_code=404, detail="Quiz not found or unauthorized")
        
        # Hash password if provided
        settings_dict = share_settings.dict()
        if settings_dict.get("password"):
            settings_dict["password"] = pwd_context.hash(settings_dict["password"])
        
        # Update quiz
        await db.quizzes.update_one(
            {"_id": ObjectId(quiz_id)},
            {"$set": {"share_settings": settings_dict}}
        )
        
        return {
            "message": "Share settings updated successfully",
            "settings": {
                "visibility": settings_dict["visibility"],
                "allow_anonymous": settings_dict["allow_anonymous"],
                "has_password": bool(settings_dict.get("password"))
            }
        }
        
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Error updating share settings: {str(e)}")

@app.put("/quizzes/{quiz_id}/timer-settings")
async def update_timer_settings(
    quiz_id: str,
    timer_settings: QuizTimerSettings,
    current_user: dict = Depends(get_current_user)
):
    """Update quiz timer settings"""
    from bson import ObjectId
    
    try:
        # Verify quiz ownership
        quiz = await db.quizzes.find_one({
            "_id": ObjectId(quiz_id),
            "created_by": current_user["email"]
        })
        
        if not quiz:
            raise HTTPException(status_code=404, detail="Quiz not found or unauthorized")
        
        settings_dict = timer_settings.dict()
        
        # Update quiz
        await db.quizzes.update_one(
            {"_id": ObjectId(quiz_id)},
            {"$set": {"timer_settings": settings_dict}}
        )
        
        return {
            "message": "Timer settings updated successfully",
            "settings": settings_dict
        }
        
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Error updating timer settings: {str(e)}")


@app.post("/public/quiz/{quiz_id}/verify-access")
async def verify_quiz_access(quiz_id: str, password: Optional[str] = None):
    """Verify access to a password-protected quiz"""
    from bson import ObjectId
    
    try:
        quiz = await db.quizzes.find_one({"_id": ObjectId(quiz_id)})
        
        if not quiz:
            raise HTTPException(status_code=404, detail="Quiz not found")
        
        share_settings = quiz.get("share_settings", {})
        
        # Check if quiz is password protected
        if share_settings.get("visibility") == "password_protected":
            if not password:
                raise HTTPException(status_code=401, detail="Password required")
            
            stored_password = share_settings.get("password")
            if not stored_password or not verify_password(password, stored_password):
                raise HTTPException(status_code=401, detail="Incorrect password")
        
        return {
            "access_granted": True,
            "message": "Access granted"
        }
        
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Error verifying access: {str(e)}")

# ==================== ANALYTICS ENDPOINTS ====================

@app.get("/analytics/quizzes")
async def get_quizzes_for_analytics(current_user: dict = Depends(get_current_user)):
    """Get all quizzes with attempt counts for analytics listing"""
    from bson import ObjectId
    
    try:
        # Get user's quizzes
        quizzes_cursor = db.quizzes.find({"created_by": current_user["email"]})
        quizzes = await quizzes_cursor.to_list(length=None)
        
        result = []
        for quiz in quizzes:
            # Count attempts for this quiz
            quiz_id = str(quiz["_id"])
            attempt_count = await db.quiz_attempts.count_documents({"quiz_id": quiz_id})
            
            result.append({
                "id": quiz_id,
                "title": quiz.get("title", "Untitled Quiz"),
                "total_questions": len(quiz.get("questions", [])),
                "total_attempts": attempt_count,
                "created_at": quiz.get("created_at", datetime.utcnow()).isoformat()
            })
        
        return {"quizzes": result}
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error fetching quizzes: {str(e)}")

@app.get("/analytics/user/overview")
async def get_user_analytics_overview(current_user: dict = Depends(get_current_user)):
    """Get overall user analytics across all quizzes"""
    from bson import ObjectId
    
    try:
        # Get all user's quizzes
        quizzes_cursor = db.quizzes.find({"created_by": current_user["email"]})
        quizzes = await quizzes_cursor.to_list(length=None)
        
        if not quizzes:
            return {
                "total_quizzes": 0,
                "total_attempts": 0,
                "total_students": 0,
                "average_score": 0,
                "quiz_performance": [],
                "recent_activity": [],
                "score_distribution": {"0-20": 0, "21-40": 0, "41-60": 0, "61-80": 0, "81-100": 0}
            }
        
        quiz_ids = [str(q["_id"]) for q in quizzes]
        
        # Get all attempts for user's quizzes
        attempts_cursor = db.quiz_attempts.find({"quiz_id": {"$in": quiz_ids}})
        attempts = await attempts_cursor.to_list(length=None)
        
        # Calculate overall stats
        total_attempts = len(attempts)
        unique_students = set(a.get("student_email") for a in attempts if a.get("student_email"))
        
        # Calculate average score
        scores = []
        for attempt in attempts:
            score = attempt.get("score", 0)
            max_score = attempt.get("max_score", 1)
            percentage = (score / max_score * 100) if max_score > 0 else 0
            scores.append(percentage)
        
        avg_score = round(sum(scores) / len(scores), 1) if scores else 0
        
        # Score distribution
        score_distribution = {"0-20": 0, "21-40": 0, "41-60": 0, "61-80": 0, "81-100": 0}
        for score in scores:
            if score <= 20:
                score_distribution["0-20"] += 1
            elif score <= 40:
                score_distribution["21-40"] += 1
            elif score <= 60:
                score_distribution["41-60"] += 1
            elif score <= 80:
                score_distribution["61-80"] += 1
            else:
                score_distribution["81-100"] += 1
        
        # Per-quiz performance
        quiz_performance = []
        for quiz in quizzes:
            quiz_id = str(quiz["_id"])
            quiz_attempts = [a for a in attempts if a.get("quiz_id") == quiz_id]
            
            if quiz_attempts:
                quiz_scores = []
                for attempt in quiz_attempts:
                    score = attempt.get("score", 0)
                    max_score = attempt.get("max_score", 1)
                    percentage = (score / max_score * 100) if max_score > 0 else 0
                    quiz_scores.append(percentage)
                
                avg_quiz_score = round(sum(quiz_scores) / len(quiz_scores), 1)
                
                quiz_performance.append({
                    "quiz_id": quiz_id,
                    "quiz_title": quiz.get("title", "Untitled"),
                    "attempts": len(quiz_attempts),
                    "avg_score": avg_quiz_score,
                    "questions": len(quiz.get("questions", []))
                })
        
        # Sort by attempts (most attempted first)
        quiz_performance.sort(key=lambda x: x["attempts"], reverse=True)
        
        # Recent activity (last 10 attempts)
        recent_attempts = sorted(attempts, key=lambda x: x.get("submitted_at", datetime.min), reverse=True)[:10]
        recent_activity = []
        for attempt in recent_attempts:
            quiz = next((q for q in quizzes if str(q["_id"]) == attempt.get("quiz_id")), None)
            if quiz:
                score = attempt.get("score", 0)
                max_score = attempt.get("max_score", 1)
                percentage = round((score / max_score * 100), 1) if max_score > 0 else 0
                
                # Get student info
                student_email = attempt.get("student_email")
                student_name = attempt.get("student_name", "Anonymous")
                
                recent_activity.append({
                    "quiz_title": quiz.get("title", "Untitled"),
                    "student_name": student_name,
                    "score": score,
                    "max_score": max_score,
                    "percentage": percentage,
                    "submitted_at": attempt.get("submitted_at").isoformat() if attempt.get("submitted_at") else None
                })
        
        return {
            "total_quizzes": len(quizzes),
            "total_attempts": total_attempts,
            "total_students": len(unique_students),
            "average_score": avg_score,
            "quiz_performance": quiz_performance,
            "recent_activity": recent_activity,
            "score_distribution": score_distribution
        }
        
    except Exception as e:
        print(f"Error in user analytics: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error fetching user analytics: {str(e)}")

@app.get("/analytics/quiz/{quiz_id}")
async def get_quiz_analytics(quiz_id: str, current_user: dict = Depends(get_current_user)):
    """Get comprehensive MCQ analytics for a specific quiz with advanced metrics"""
    from bson import ObjectId
    from collections import defaultdict
    
    try:
        # Verify quiz ownership
        quiz = await db.quizzes.find_one({
            "_id": ObjectId(quiz_id),
            "created_by": current_user["email"]
        })
        
        if not quiz:
            raise HTTPException(status_code=404, detail="Quiz not found")
        
        # Get all attempts
        attempts = await db.quiz_attempts.find({"quiz_id": quiz_id}).to_list(length=None)
        
        if not attempts or len(attempts) < 10:
            return {
                "quiz_id": quiz_id,
                "quiz_title": quiz.get("title", "Untitled Quiz"),
                "total_questions": len(quiz.get("questions", [])),
                "total_attempts": len(attempts),
                "total_students": len(set(a.get("student_email", "") for a in attempts)) if attempts else 0,
                "average_score": 0,
                "average_percentage": 0,
                "average_time_per_question": 0,
                "score_distribution": {},
                "insufficient_data": True,
                "message": "Insufficient data to analyze. At least 10 attempts required.",
                "questions": []
            }
        
        # Calculate quiz-level stats
        total_students = len(set(a.get("student_email", "") for a in attempts))
        scores = [a.get("score", 0) for a in attempts]
        percentages = [a.get("percentage", 0) for a in attempts]
        avg_score = sum(scores) / len(scores) if scores else 0
        avg_percentage = sum(percentages) / len(percentages) if percentages else 0
        
        # Calculate average time per question
        total_time = sum(a.get("time_taken", 0) for a in attempts)
        total_questions = len(quiz.get("questions", []))
        avg_time_per_question = round(total_time / (len(attempts) * total_questions), 1) if total_questions > 0 else 0
        
        # Score distribution (group by ranges)
        score_distribution = defaultdict(int)
        for score in scores:
            if score <= 3:
                score_distribution["0-3"] += 1
            elif score <= 6:
                score_distribution["4-6"] += 1
            elif score <= 9:
                score_distribution["7-9"] += 1
            else:
                score_distribution["10+"] += 1
        
        # Sort attempts by score for discrimination index calculation
        sorted_attempts = sorted(attempts, key=lambda x: x.get("score", 0), reverse=True)
        top_25_count = max(1, len(sorted_attempts) // 4)
        top_25_percent = sorted_attempts[:top_25_count]
        bottom_25_percent = sorted_attempts[-top_25_count:]
        
        # Analyze each question
        questions = quiz.get("questions", [])
        question_analytics = []
        
        for q_idx, question in enumerate(questions):
            # Count responses for this question
            correct_count = 0
            incorrect_count = 0
            skipped_count = 0
            question_times = []
            
            # Track option selections
            options = question.get("options", [])
            option_stats = defaultdict(int)
            
            # Discrimination index calculation
            top_25_correct = 0
            bottom_25_correct = 0
            
            for attempt in attempts:
                # Check if correct
                if q_idx in attempt.get("correct_answers", []):
                    correct_count += 1
                elif q_idx in attempt.get("incorrect_answers", []):
                    incorrect_count += 1
                elif q_idx in attempt.get("unanswered", []):
                    skipped_count += 1
                
                # Track selected options
                user_answers = attempt.get("answers", {}).get(str(q_idx), [])
                for opt_idx in user_answers:
                    option_stats[opt_idx] += 1
                
                # Track time spent on question (if available)
                time_per_q = attempt.get("time_per_question", {})
                question_time = time_per_q.get(str(q_idx), 0) if isinstance(time_per_q, dict) else 0
                if question_time > 0:
                    question_times.append(question_time)
            
            # Validate data integrity
            total_responses = len(attempts)
            actual_total = correct_count + incorrect_count + skipped_count
            if actual_total != total_responses:
                print(f"⚠️ Warning: Q{q_idx+1} data mismatch - Total: {total_responses}, Sum: {actual_total}")
            
            # Calculate discrimination index with fallback
            discrimination_index = None
            if top_25_count >= 5 and len(attempts) >= 20:  # Need sufficient data
                for attempt in top_25_percent:
                    if q_idx in attempt.get("correct_answers", []):
                        top_25_correct += 1
                
                for attempt in bottom_25_percent:
                    if q_idx in attempt.get("correct_answers", []):
                        bottom_25_correct += 1
                
                discrimination_index = round(
                    (top_25_correct / top_25_count - bottom_25_correct / top_25_count), 1
                )
            
            # Ensure precise rounding to 1 decimal place
            correct_percentage = round((correct_count / total_responses * 100) if total_responses > 0 else 0, 1)
            incorrect_percentage = round((incorrect_count / total_responses * 100) if total_responses > 0 else 0, 1)
            skip_rate = round((skipped_count / total_responses * 100) if total_responses > 0 else 0, 1)
            avg_time_spent = round(sum(question_times) / len(question_times), 1) if question_times else None
            
            # Determine correct option indices
            correct_option_indices = [i for i, opt in enumerate(options) if opt.get("is_correct", False)]
            
            # Build option breakdown
            option_breakdown = []
            most_chosen_option = None
            max_selections = 0
            
            for opt_idx, option in enumerate(options):
                selection_count = option_stats.get(opt_idx, 0)
                selection_percentage = round((selection_count / total_responses) * 100, 1) if total_responses > 0 else 0
                
                if selection_count > max_selections:
                    max_selections = selection_count
                    most_chosen_option = chr(65 + opt_idx)
                
                option_breakdown.append({
                    "option_index": opt_idx,
                    "option_letter": chr(65 + opt_idx),  # A, B, C, D
                    "option_text": option.get("text", ""),
                    "is_correct": opt_idx in correct_option_indices,
                    "selected_count": selection_count,
                    "selection_percentage": selection_percentage
                })
            
            # Flag problematic questions with detailed logic
            is_problematic = False
            flags = []
            flag_emoji = "good"
            flag_icon = "✅"
            flag_tooltip = "Good question - performing well"
            
            if correct_percentage < 40:
                flags.append("Too Difficult")
                is_problematic = True
                flag_emoji = "warning"
                flag_icon = "⚠️"
                flag_tooltip = "Too Difficult: Less than 40% correct"
            
            if discrimination_index is not None and discrimination_index < 0.2:
                if not flags:
                    flags.append("Needs Review (Low Discrimination)")
                    flag_tooltip = "Needs Review: Low discrimination index"
                is_problematic = True
                flag_emoji = "warning"
                flag_icon = "⚠️"
            
            if skip_rate > 20:
                flags.append("Unclear Question")
                is_problematic = True
                flag_emoji = "unclear"
                flag_icon = "❗"
                flag_tooltip = "Unclear: High skip rate (>20%)"
            
            question_analytics.append({
                "question_index": q_idx,
                "question_number": q_idx + 1,
                "question_text": question.get("question", "") or "No question text",
                "attempts": total_responses,
                "correct_count": correct_count,
                "incorrect_count": incorrect_count,
                "skipped_count": skipped_count,
                "correct_percentage": correct_percentage,
                "incorrect_percentage": incorrect_percentage,
                "skip_rate": skip_rate,
                "average_time_spent": avg_time_spent if avg_time_spent is not None else 0,
                "discrimination_index": discrimination_index if discrimination_index is not None else "N/A",
                "most_chosen_option": most_chosen_option or "—",
                "is_problematic": is_problematic,
                "flag_emoji": flag_emoji,
                "flag_icon": flag_icon,
                "flag_tooltip": flag_tooltip,
                "flags": flags,
                "options": option_breakdown
            })
        
        return {
            "quiz_id": quiz_id,
            "quiz_title": quiz.get("title", "Untitled Quiz"),
            "total_questions": len(questions),
            "total_attempts": len(attempts),
            "total_students": total_students,
            "average_score": round(avg_score, 1),
            "average_percentage": round(avg_percentage, 1),
            "average_time_per_question": avg_time_per_question,
            "score_distribution": dict(score_distribution),
            "questions": question_analytics
        }
        
    except HTTPException:
        raise
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Error generating analytics: {str(e)}")

@app.get("/analytics/quiz/{quiz_id}/export")
async def export_quiz_analytics(quiz_id: str, current_user: dict = Depends(get_current_user)):
    """Export quiz analytics as CSV with timestamp"""
    from bson import ObjectId
    import io
    import csv
    
    try:
        # Get analytics data
        analytics = await get_quiz_analytics(quiz_id, current_user)
        
        # Check for insufficient data
        if analytics.get('insufficient_data'):
            raise HTTPException(status_code=400, detail="Insufficient data to export. At least 10 attempts required.")
        
        output = io.StringIO()
        writer = csv.writer(output)
        
        # Write header with timestamp
        export_time = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        writer.writerow([f"Quiz Analytics: {analytics['quiz_title']}"])
        writer.writerow([f"Exported: {export_time}"])
        writer.writerow([])
        
        # Write summary
        writer.writerow(['Total Questions', analytics['total_questions']])
        writer.writerow(['Total Attempts', analytics['total_attempts']])
        writer.writerow(['Total Students', analytics['total_students']])
        writer.writerow(['Average Score', f"{analytics['average_score']}/{analytics['total_questions']}"])
        writer.writerow(['Average Percentage', f"{analytics['average_percentage']}%"])
        writer.writerow(['Average Time per Question', f"{analytics.get('average_time_per_question', 0)}s"])
        writer.writerow([])
        writer.writerow(['Score Distribution'])
        for range_label, count in analytics.get('score_distribution', {}).items():
            writer.writerow([f"  {range_label}", count])
        writer.writerow([])
        
        # Write question-level data with consistent headers
        writer.writerow(['Question ID', 'Question Text', 'Attempts', 'Correct %', 'Skip %', 'Discrimination', 'Flag', 'Most Chosen Option'])
        
        for q in analytics['questions']:
            disc_value = q.get('discrimination_index', 'N/A')
            disc_display = f"{disc_value}" if disc_value != 'N/A' else 'N/A'
            
            flag_text = q.get('flag_icon', '✅') + ' ' + ', '.join(q.get('flags', [])) if q.get('flags') else '✅ Good'
            
            writer.writerow([
                f"Q{q['question_number']}",
                q.get('question_text', 'No text')[:100],
                q.get('attempts', 0),
                f"{q.get('correct_percentage', 0)}%",
                f"{q.get('skip_rate', 0)}%",
                disc_display,
                flag_text,
                q.get('most_chosen_option', '—')
            ])
        
        writer.writerow([])
        writer.writerow(['Option-Level Analysis'])
        writer.writerow([])
        
        # Write option-level data
        for q in analytics['questions']:
            writer.writerow([f"Q{q['question_number']}: {q['question_text'][:80]}"])
            writer.writerow(['Option', 'Text', 'Correct', 'Selections', 'Selection %'])
            
            for opt in q['options']:
                writer.writerow([
                    opt['option_letter'],
                    opt['option_text'][:100],
                    'Yes' if opt['is_correct'] else 'No',
                    opt['selection_count'],
                    f"{opt['selection_percentage']}%"
                ])
            
            writer.writerow([])
        
        output.seek(0)
        
        # Generate filename with timestamp
        timestamp = datetime.now().strftime("%Y-%m-%d_%H-%M")
        quiz_title_safe = analytics['quiz_title'].replace(' ', '_').replace('/', '-')[:30]
        filename = f"{quiz_title_safe}_Analytics_{timestamp}.csv"
        
        return StreamingResponse(
            iter([output.getvalue()]),
            media_type="text/csv",
            headers={
                "Content-Disposition": f"attachment; filename={filename}"
            }
        )
        
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error exporting analytics: {str(e)}")

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
